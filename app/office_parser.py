"""Office/Web Parser: ekstraksi teks dari .docx/.pptx/.html + smart chunking.

Konsisten dengan md_parser & pdf_parser: teks dipecah per heading
(docx: style "Heading N", html: tag h1-h6, pptx: judul slide), lalu segmen
panjang dipecah dengan LangChain ``RecursiveCharacterTextSplitter``
(parameter sama persis) supaya chunk konsisten lintas format.

Metadata yang dihasilkan: ``{"source", "page", "heading", "chunk_index"}``.
page selalu 1 untuk docx/html; untuk pptx = nomor slide (seperti PDF
per halaman).
"""

from __future__ import annotations

import logging
from html.parser import HTMLParser
from pathlib import Path

from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from app.pdf_parser import DEFAULT_CHUNK_OVERLAP, DEFAULT_CHUNK_SIZE, Chunk

logger = logging.getLogger(__name__)

NO_HEADING_LABEL = "Intro"

# Tag yang isinya tidak informatif untuk RAG (script, nav, dsb).
_SKIP_TAGS = {"script", "style", "noscript", "template", "head", "iframe", "svg"}
_HEADING_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6"}
# Tag yang memulai paragraf/baris baru dalam alur teks HTML.
_BLOCK_TAGS = {
    "p", "div", "li", "ul", "ol", "tr", "br", "section", "article",
    "blockquote", "table", "pre", "hr",
}


def _chunk_segments(
    segments: list[tuple[str | None, str, int]],
    source: str,
    chunk_size: int,
    chunk_overlap: int,
) -> list[Chunk]:
    """Ubah segmen (heading, teks, page) menjadi chunk ber-metadata.

    Parameter chunking sama persis dengan pdf_parser & md_parser.
    Segmen yang isinya hanya judul (tanpa isi) di-skip.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks: list[Chunk] = []
    for heading, segment, page in segments:
        if segment == heading:
            continue  # segmen berisi judul saja, tanpa isi
        parts = [segment] if len(segment) <= chunk_size else splitter.split_text(segment)
        for part in parts:
            chunks.append(
                Chunk(
                    text=part,
                    metadata={"page": page, "heading": heading or NO_HEADING_LABEL},
                )
            )

    for idx, chunk in enumerate(chunks):
        chunk.metadata["source"] = source
        chunk.metadata["chunk_index"] = idx
    return chunks


def parse_docx(
    path: str | Path,
    source: str | None = None,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[Chunk]:
    """Baca .docx, pecah per paragraf ber-style Heading, lalu chunk.

    - Paragraf ber-style "Heading N" memulai segmen baru (judul ikut
      disertakan sebagai baris pertama isi, sama seperti markdown).
    - Isi tabel ikut diekstrak: tiap baris jadi satu baris teks dengan
      sel dipisah " | " (urutan asli dokumen dipertahankan).
    - Teks sebelum heading pertama ber-heading "Intro".

    Args:
        path: Lokasi file .docx.
        source: Nama sumber (default: nama file).
        chunk_size: Ukuran maksimum chunk (karakter).
        chunk_overlap: Jumlah karakter overlap antar chunk.

    Returns:
        List of :class:`Chunk` (page selalu 1).
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File tidak ditemukan: {path}")
    source = source or path.name

    doc = Document(str(path))
    segments: list[tuple[str | None, str, int]] = []
    current_heading: str | None = None
    current_parts: list[str] = []

    def flush() -> None:
        nonlocal current_heading, current_parts
        body = "\n".join(current_parts).strip()
        if body:
            segments.append((current_heading, body, 1))
        current_heading = None
        current_parts = []

    for block in doc.iter_inner_content():
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style_name = (block.style.name if block.style else "") or ""
            if style_name.lower().startswith("heading"):
                flush()
                current_heading = text
                current_parts = [text]
            else:
                current_parts.append(text)
        elif isinstance(block, Table):
            for row in block.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells).strip()
                if row_text:
                    current_parts.append(row_text)
    flush()

    return _chunk_segments(segments, source, chunk_size, chunk_overlap)


def parse_pptx(
    path: str | Path,
    source: str | None = None,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[Chunk]:
    """Baca .pptx: tiap slide jadi satu segmen, lalu chunk.

    - Judul slide (title placeholder) dijadikan heading; fallback
      "Slide N". Teks shape lain digabung sebagai isi.
    - page = nomor slide (1-based), konsisten dengan PDF per halaman.
    - Slide tanpa teks sama sekali di-skip.

    Args:
        path: Lokasi file .pptx.
        source: Nama sumber (default: nama file).
        chunk_size: Ukuran maksimum chunk (karakter).
        chunk_overlap: Jumlah karakter overlap antar chunk.

    Returns:
        List of :class:`Chunk`.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File tidak ditemukan: {path}")
    source = source or path.name

    prs = Presentation(str(path))
    segments: list[tuple[str | None, str, int]] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = ""
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()

        parts: list[str] = []
        for shape in slide.shapes:
            if shape is title_shape or not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

        body = "\n".join(parts).strip()
        if not body and not title_text:
            continue  # slide kosong
        heading = title_text or f"Slide {slide_no}"
        segment = "\n".join(part for part in (title_text, body) if part)
        segments.append((heading, segment, slide_no))

    return _chunk_segments(segments, source, chunk_size, chunk_overlap)


class _HtmlToSegments(HTMLParser):
    """Ubah HTML menjadi segmen (heading, teks) berbasis h1-h6.

    - Tag dalam _SKIP_TAGS (script, style, dll.) ikut isinya dilewati.
    - h1-h6 memulai segmen baru; teks heading ikut sebagai baris pertama.
    - Tag block memulai baris baru; pre mempertahankan whitespace asli.
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.segments: list[tuple[str | None, str]] = []
        self._current_heading: str | None = None
        self._current_parts: list[str] = []
        self._skip_depth = 0
        self._heading_parts: list[str] | None = None
        self._need_break = False
        self._in_pre = False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if tag in _SKIP_TAGS:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return
        if tag in _HEADING_TAGS:
            self._flush_segment()
            self._heading_parts = []
            self._need_break = False
            return
        if tag in _BLOCK_TAGS:
            self._need_break = True
            if tag == "pre":
                self._in_pre = True

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in _SKIP_TAGS:
            if self._skip_depth:
                self._skip_depth -= 1
            return
        if self._skip_depth:
            return
        if tag in _HEADING_TAGS:
            if self._heading_parts is not None:
                heading = " ".join(" ".join(self._heading_parts).split())
                if heading:
                    self._current_heading = heading
                    self._current_parts.append(heading)
            self._heading_parts = None
            self._need_break = True
            return
        if tag == "pre":
            self._in_pre = False
            self._need_break = True
        elif tag in _BLOCK_TAGS:
            self._need_break = True

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        if self._heading_parts is not None:
            self._heading_parts.append(data)
            return
        if not data.strip():
            return
        if self._need_break:
            if self._current_parts:
                self._current_parts.append("")
            self._need_break = False
        if self._in_pre:
            self._current_parts.append(data)
        else:
            text = " ".join(data.split())
            if self._current_parts and self._current_parts[-1]:
                self._current_parts[-1] += " " + text
            else:
                self._current_parts.append(text)

    def _flush_segment(self) -> None:
        body = "\n".join(self._current_parts).strip()
        if body:
            self.segments.append((self._current_heading, body))
        self._current_heading = None
        self._current_parts = []
        self._need_break = False

    def finish(self) -> None:
        """Flush segmen terakhir setelah seluruh HTML selesai di-feed."""
        self._flush_segment()


def parse_html(
    path: str | Path,
    source: str | None = None,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[Chunk]:
    """Baca .html/.htm, pecah per heading h1-h6, lalu chunk.

    - Script/style/head/iframe/svg beserta isinya dilewati.
    - Tag block (p, li, div, dll.) jadi pemisah baris; teks dalam satu
      paragraf digabung dengan spasi.
    - pre mempertahankan whitespace asli (berguna untuk cuplikan kode).
    - Teks sebelum h1 pertama ber-heading "Intro".

    Args:
        path: Lokasi file .html/.htm.
        source: Nama sumber (default: nama file).
        chunk_size: Ukuran maksimum chunk (karakter).
        chunk_overlap: Jumlah karakter overlap antar chunk.

    Returns:
        List of :class:`Chunk` (page selalu 1).
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File tidak ditemukan: {path}")
    source = source or path.name

    text = path.read_text(encoding="utf-8-sig", errors="replace")
    parser = _HtmlToSegments()
    parser.feed(text)
    parser.finish()

    segments = [(heading, segment, 1) for heading, segment in parser.segments]
    return _chunk_segments(segments, source, chunk_size, chunk_overlap)
