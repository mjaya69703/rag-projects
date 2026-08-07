"""Test fitur: office_parser — parse_docx, parse_pptx, parse_html.

Jalankan: python tests/test_office_parser.py  atau  pytest tests/test_office_parser.py -v
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.office_parser import NO_HEADING_LABEL, parse_docx, parse_html, parse_pptx
from app.watch_folder import parse_any

_LONG_BODY = (
    "VLAN adalah metode mempartisi satu jaringan fisik menjadi beberapa jaringan "
    "logis. Dengan VLAN, broadcast domain diperkecil sehingga lalu lintas jaringan "
    "lebih efisien. Setiap VLAN memiliki identitas nomor antara 1 hingga 4094, dan "
    "komunikasi antar VLAN memerlukan perangkat layer 3. Konfigurasi dilakukan pada "
    "port switch dengan mode access atau trunk. Port access melewatkan trafik satu "
    "VLAN saja, sedangkan port trunk melewatkan banyak VLAN menggunakan tag 802.1Q. "
    "Mode trunk dipakai ketika menghubungkan switch ke switch lain atau ke router. "
    "Perencanaan yang matang mencegah masalah routing sejak awal sebelum berdampak "
    "ke pengguna akhir jaringan yang sedang beroperasi sehari-hari."
)


def _make_docx(tmp_path: Path) -> Path:
    doc = Document()
    doc.add_paragraph("Catatan awal tanpa heading. Konteks pembuka dokumen.")
    doc.add_heading("Panduan Jaringan Komputer", level=1)
    doc.add_paragraph("Pengantar jaringan: kumpulan perangkat yang saling terhubung.")
    doc.add_heading("VLAN", level=2)
    doc.add_paragraph(_LONG_BODY)
    doc.add_heading("Tabel Alamat", level=3)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Nama"
    table.cell(0, 1).text = "IP Address"
    table.cell(1, 0).text = "Router A"
    table.cell(1, 1).text = "192.168.1.1"
    path = tmp_path / "materi.docx"
    doc.save(str(path))
    return path


def _make_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Pendahuluan"
    box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    box.text_frame.text = _LONG_BODY
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
    box2.text_frame.text = "Isi slide kedua tanpa judul."
    path = tmp_path / "materi.pptx"
    prs.save(str(path))
    return path


def _make_html(tmp_path: Path) -> Path:
    path = tmp_path / "halaman.html"
    path.write_text(
        """<!DOCTYPE html>
<html>
<head><title>Judul Tab</title><meta charset="utf-8"></head>
<body>
<script>var rahasia = "TIDAK BOLEH MUNCUL";</script>
<h1>Panduan Jaringan Komputer</h1>
<p>Pengantar jaringan: <b>perangkat</b> yang saling terhubung untuk berbagi data.</p>
<h2>VLAN</h2>
<p>VLAN mempartisi satu jaringan fisik menjadi beberapa jaringan logis.</p>
<pre>
def contoh():
    return True
</pre>
</body>
</html>
""",
        encoding="utf-8",
    )
    return path


def test_docx_headings_detected_with_intro(tmp_path: Path) -> None:
    chunks = parse_docx(_make_docx(tmp_path))
    headings = {c.metadata["heading"] for c in chunks}
    assert "Panduan Jaringan Komputer" in headings
    assert "VLAN" in headings
    assert "Tabel Alamat" in headings
    intro = [c for c in chunks if c.metadata["heading"] == NO_HEADING_LABEL]
    assert intro, "teks sebelum heading pertama harus ber-heading 'Intro'"
    assert "Catatan awal tanpa heading" in intro[0].text


def test_docx_table_content_extracted(tmp_path: Path) -> None:
    chunks = parse_docx(_make_docx(tmp_path))
    all_text = "\n".join(c.text for c in chunks)
    assert "192.168.1.1" in all_text
    assert "Router A" in all_text


def test_docx_long_segment_split_and_length(tmp_path: Path) -> None:
    chunks = parse_docx(_make_docx(tmp_path))
    vlan = [c for c in chunks if c.metadata["heading"] == "VLAN"]
    assert len(vlan) >= 2, "segmen panjang harus dipecah jadi beberapa chunk"
    for chunk in chunks:
        assert len(chunk.text) <= 550, f"chunk terlalu panjang: {len(chunk.text)}"
        assert chunk.text.strip(), "chunk tidak boleh kosong"


def test_docx_metadata_consistent(tmp_path: Path) -> None:
    chunks = parse_docx(_make_docx(tmp_path))
    assert chunks, "parse_docx() tidak menghasilkan chunk"
    for i, chunk in enumerate(chunks):
        meta = chunk.metadata
        assert meta["source"] == "materi.docx"
        assert meta["page"] == 1
        assert meta["chunk_index"] == i
        assert meta["heading"], f"heading kosong di chunk {i}"
        assert set(meta) == {"source", "page", "heading", "chunk_index"}


def test_pptx_per_slide_page_and_heading(tmp_path: Path) -> None:
    chunks = parse_pptx(_make_pptx(tmp_path))
    pages = [c.metadata["page"] for c in chunks]
    assert pages and min(pages) == 1 and max(pages) == 2, pages
    headings = {c.metadata["heading"] for c in chunks}
    assert "Pendahuluan" in headings
    assert "Slide 2" in headings  # slide tanpa judul → fallback "Slide N"
    assert all(c.metadata["source"] == "materi.pptx" for c in chunks)
    assert all(c.metadata["chunk_index"] == i for i, c in enumerate(chunks))


def test_pptx_long_slide_split(tmp_path: Path) -> None:
    chunks = parse_pptx(_make_pptx(tmp_path))
    first = [c for c in chunks if c.metadata["heading"] == "Pendahuluan"]
    assert len(first) >= 2, "slide panjang harus dipecah jadi beberapa chunk"


def test_html_headings_skip_script_and_keep_pre(tmp_path: Path) -> None:
    chunks = parse_html(_make_html(tmp_path))
    headings = {c.metadata["heading"] for c in chunks}
    assert "Panduan Jaringan Komputer" in headings
    assert "VLAN" in headings
    all_text = "\n".join(c.text for c in chunks)
    assert "TIDAK BOLEH MUNCUL" not in all_text  # isi <script> di-skip
    assert "Judul Tab" not in all_text  # isi <head> di-skip
    assert "def contoh():" in all_text  # isi <pre> dipertahankan
    assert all(c.metadata["page"] == 1 for c in chunks)
    assert all(c.metadata["source"] == "halaman.html" for c in chunks)


def test_parse_any_dispatches_new_formats(tmp_path: Path) -> None:
    assert parse_any(_make_docx(tmp_path))
    assert parse_any(_make_pptx(tmp_path))
    assert parse_any(_make_html(tmp_path))


def test_missing_file_raises(tmp_path: Path) -> None:
    missing = tmp_path / "tidak-ada.docx"
    with pytest.raises(FileNotFoundError):
        parse_docx(missing)
    with pytest.raises(FileNotFoundError):
        parse_pptx(tmp_path / "tidak-ada.pptx")
    with pytest.raises(FileNotFoundError):
        parse_html(tmp_path / "tidak-ada.html")


def test_source_override(tmp_path: Path) -> None:
    chunks = parse_docx(_make_docx(tmp_path), source="nama-kustom")
    assert all(c.metadata["source"] == "nama-kustom" for c in chunks)
